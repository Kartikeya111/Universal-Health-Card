from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Unveiling [Your Name]"

# Slide 2: Main Content Slide
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# Background Image
img_path = 'cityscape.jpeg'  # Add path to your background image
left = top = Inches(0)
pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

# Main Text Box
left = Inches(1)
top = Inches(1)
width = prs.slide_width - Inches(2)
height = prs.slide_height - Inches(2)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame

# Main Text Content
main_text = text_frame.add_paragraph()
main_text.text = "Name: [Your Name]"
main_text.level = 0
main_text.space_after = Pt(16)

main_text = text_frame.add_paragraph()
main_text.text = "- Welcome! Let's dive into the world of a problem solver and data enthusiast."
main_text.level = 0
main_text.space_after = Pt(12)

# Add other content similarly

# Define animations
def add_animation(shape, idx, animation_type):
    effect = shape.text_frame.paragraphs[idx].text = "Passions:"
    animation = shape.text_frame.paragraphs[idx].runs[0].add_animation(animation_type)
    return effect, animation

# Add animations
textbox, idx = slide.shapes[1], 0
effect, animation = add_animation(textbox, idx, MSO_ANIMATION_TYPE.FADE)

# Save presentation
prs.save('presentation.pptx')
